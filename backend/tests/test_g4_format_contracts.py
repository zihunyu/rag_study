from __future__ import annotations

import hashlib
import wave
from pathlib import Path

import openpyxl
from docx import Document
from pptx import Presentation
from pypdf import PdfWriter
from ragkb.document_processing.parsers import ParserRouter
from ragkb.domain.validation import DocumentQualityReport, QualityDisposition
from ragkb.engineering_security.file_validation import UploadFileValidator


def _synthetic_files(tmp_path: Path) -> dict[str, Path]:
    files: dict[str, Path] = {}
    for source_format, name, content in (
        ("txt", "fixture.txt", b"synthetic text\nline two"),
        ("markdown", "fixture.md", b"# Synthetic\nbody"),
        ("html", "fixture.html", b"<h1>Synthetic</h1><p>body</p>"),
        ("image", "fixture.png", b"\x89PNG\r\n\x1a\nsynthetic"),
        ("doc", "fixture.doc", b"synthetic legacy word fixture"),
        ("ppt", "fixture.ppt", b"synthetic legacy slides fixture"),
        ("csv", "fixture.csv", b"name,value\nalpha,1\n"),
    ):
        path = tmp_path / name
        path.write_bytes(content)
        files[source_format] = path

    pdf = tmp_path / "scanned.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    with pdf.open("wb") as handle:
        writer.write(handle)
    files["pdf"] = pdf

    docx = tmp_path / "fixture.docx"
    document = Document()
    document.add_paragraph("synthetic docx")
    document.save(docx)
    files["docx"] = docx

    pptx = tmp_path / "fixture.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "synthetic pptx"
    presentation.save(pptx)
    files["pptx"] = pptx

    xlsx = tmp_path / "fixture.xlsx"
    workbook = openpyxl.Workbook()
    workbook.active.append(["name", "value"])
    workbook.active.append(["alpha", 1])
    workbook.save(xlsx)
    workbook.close()
    files["xlsx"] = xlsx

    wav = tmp_path / "fixture.wav"
    with wave.open(str(wav), "wb") as audio:
        audio.setnchannels(1)
        audio.setsampwidth(2)
        audio.setframerate(8000)
        audio.writeframes(b"\x00\x00" * 800)
    files["audio"] = wav
    return files


def test_all_g4_local_format_routes_emit_locator_and_quality_contracts(tmp_path: Path) -> None:
    router = ParserRouter()
    documents = {
        source_format: router.parse(source_format, path, f"version-{source_format}")
        for source_format, path in _synthetic_files(tmp_path).items()
    }
    for extension, content in (
        ("mp3", b"ID3synthetic-mp3"),
        ("m4a", b"\x00\x00\x00\x18ftypM4Asynthetic"),
    ):
        audio_path = tmp_path / f"fixture.{extension}"
        audio_path.write_bytes(content)
        documents[f"audio_{extension}"] = router.parse(
            "audio", audio_path, f"version-audio-{extension}"
        )

    assert {
        "pdf",
        "image",
        "doc",
        "docx",
        "ppt",
        "pptx",
        "html",
        "markdown",
        "txt",
        "xlsx",
        "csv",
        "audio",
    }.issubset(documents)
    reports = [DocumentQualityReport.from_document(document) for document in documents.values()]
    assert all(report.locator_coverage == 1.0 for report in reports)
    assert all(report.real_acceptance is False for report in reports)
    assert documents["audio_mp3"].nodes[0].locator.start_time == 0.0
    assert documents["audio_m4a"].nodes[0].locator.end_time == 1.0
    blocked_formats = {
        report.source_format
        for report in reports
        if report.disposition is QualityDisposition.BLOCKED_REAL_VALIDATION
    }
    assert {"pdf_scanned", "image", "doc", "ppt", "audio"}.issubset(blocked_formats)


def test_wav_mp3_m4a_upload_magic_contracts_are_local_only(tmp_path: Path) -> None:
    validator = UploadFileValidator(max_size_bytes=1024 * 1024)
    fixtures = {
        "fixture.wav": (b"RIFF\x10\x00\x00\x00WAVEsynthetic", "audio/wav"),
        "fixture.mp3": (b"ID3synthetic", "audio/mpeg"),
        "fixture.m4a": (b"\x00\x00\x00\x18ftypM4Asynthetic", "audio/mp4"),
    }
    for filename, (content, mime) in fixtures.items():
        path = tmp_path / filename
        path.write_bytes(content)
        detected = validator.inspect(
            path,
            filename=filename,
            expected_size=len(content),
            expected_sha256=hashlib.sha256(content).hexdigest(),
            declared_mime=mime,
        )
        assert detected.source_format == "audio"
