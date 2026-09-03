from __future__ import annotations

import json
from pathlib import Path

import jsonschema
import openpyxl
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
from ragkb.document_processing.parsers import ParserRouter
from ragkb.domain.ids import new_uuid7


def _assert_contract(document, repository_root: Path) -> None:
    schema = json.loads(
        (
            repository_root
            / "backend/src/ragkb/contracts/schemas/canonical-document-v1.schema.json"
        ).read_text(encoding="utf-8")
    )
    jsonschema.Draft202012Validator(schema).validate(document.to_dict())
    assert document.real_acceptance is False
    assert document.nodes


def test_text_markdown_and_html_routes(tmp_path: Path) -> None:
    router = ParserRouter()
    markdown = tmp_path / "sample.md"
    markdown.write_text("# Heading\nParagraph", encoding="utf-8")
    html = tmp_path / "sample.html"
    html.write_text(
        "<html><script>ignore()</script><body><h1>Title</h1><p>Visible</p></body></html>",
        encoding="utf-8",
    )

    markdown_doc = router.parse("markdown", markdown, new_uuid7())
    html_doc = router.parse("html", html, new_uuid7())

    assert len(markdown_doc.nodes) == 2
    assert [node.original_text for node in html_doc.nodes] == ["Title", "Visible"]
    assert all(node.locator.char_range is not None for node in markdown_doc.nodes)


def test_xlsx_and_csv_preserve_sheet_row_and_formula(tmp_path: Path) -> None:
    router = ParserRouter()
    xlsx = tmp_path / "sample.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Budget"
    sheet.append(["Item", "Amount"])
    sheet.append(["Total", "=SUM(B3:B3)"])
    workbook.save(xlsx)
    csv_path = tmp_path / "sample.csv"
    csv_path.write_text("name,value\nalpha,1\n", encoding="utf-8")

    xlsx_doc = router.parse("xlsx", xlsx, new_uuid7())
    csv_doc = router.parse("csv", csv_path, new_uuid7())

    assert xlsx_doc.nodes[0].locator.sheet == "Budget"
    assert xlsx_doc.nodes[0].locator.cell_range == "A1:B1"
    assert "=SUM(B3:B3)" in xlsx_doc.nodes[1].original_text
    assert csv_doc.nodes[1].locator.sheet == "csv"
    assert csv_doc.tables[0]["non_empty_rows"] == 2


def test_spreadsheet_and_csv_keep_source_row_column_and_record_sequence(tmp_path: Path) -> None:
    csv_path = tmp_path / "quoted.csv"
    csv_path.write_text('a,b,c\nfirst,"line-one\nline-two",last\nnext,mid,end\n', encoding="utf-8")
    workbook = openpyxl.Workbook()
    first = workbook.active
    first.title = "First"
    first.append(["left", "right", "tail"])
    second = workbook.create_sheet("Second")
    second.append(["other", "record"])
    xlsx_path = tmp_path / "sequence.xlsx"
    workbook.save(xlsx_path)

    router = ParserRouter()
    csv_doc = router.parse("csv", csv_path, new_uuid7())
    xlsx_doc = router.parse("xlsx", xlsx_path, new_uuid7())

    assert [node.locator.row for node in csv_doc.nodes] == [1, 2, 3]
    assert csv_doc.nodes[1].metadata["column_addresses"] == ["A", "B", "C"]
    assert "line-one\nline-two" in csv_doc.nodes[1].display_text.replace("\r\n", "\n")
    assert [node.locator.sheet for node in xlsx_doc.nodes] == ["First", "Second"]
    assert xlsx_doc.nodes[0].locator.cell_range == "A1:C1"


def test_docx_and_pptx_routes_have_traceable_locators(tmp_path: Path) -> None:
    router = ParserRouter()
    docx_path = tmp_path / "sample.docx"
    docx = Document()
    docx.add_heading("Policy", level=1)
    docx.add_paragraph("Warranty is three years.")
    docx.save(docx_path)
    pptx_path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Roadmap"
    presentation.save(pptx_path)

    docx_result = router.parse("docx", docx_path, new_uuid7())
    pptx_result = router.parse("pptx", pptx_path, new_uuid7())

    assert docx_result.nodes[0].locator.char_range is not None
    assert "page_mapping_unavailable_in_native_docx_route" in docx_result.quality_issues
    assert pptx_result.nodes[0].locator.slide == 1


def test_pptx_spatial_cards_and_table_preserve_container_reading_order(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    left_top = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.5))
    left_top.text_frame.text = "card-left-top"
    left_bottom = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(2), Inches(0.5))
    left_bottom.text_frame.text = "card-left-bottom"
    right_top = slide.shapes.add_textbox(Inches(4.5), Inches(0.5), Inches(2), Inches(0.5))
    right_top.text_frame.text = "card-right-top"
    table = slide.shapes.add_table(2, 2, Inches(4.5), Inches(2), Inches(2), Inches(1)).table
    table.cell(0, 0).text = "table-a"
    table.cell(0, 1).text = "table-b"
    table.cell(1, 0).text = "table-c"
    table.cell(1, 1).text = "table-d"
    path = tmp_path / "cards.pptx"
    presentation.save(path)

    parsed = ParserRouter().parse("pptx", path, new_uuid7())

    assert [node.display_text for node in parsed.nodes] == [
        "card-left-top",
        "card-left-bottom",
        "card-right-top",
        "table-a | table-b\ntable-c | table-d",
    ]


def test_blank_pdf_and_image_routes_use_offline_stub_without_claiming_real_support(
    tmp_path: Path,
) -> None:
    router = ParserRouter()
    pdf = tmp_path / "blank.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    with pdf.open("wb") as handle:
        writer.write(handle)
    image = tmp_path / "image.png"
    image.write_bytes(b"\x89PNG\r\n\x1a\n")

    pdf_result = router.parse("pdf", pdf, new_uuid7())
    image_result = router.parse("image", image, new_uuid7())

    assert pdf_result.source_format == "pdf_scanned"
    assert image_result.source_format == "image"
    assert "offline_ocr_stub_real_effect_blocked" in pdf_result.quality_issues
    assert "offline_ocr_stub_real_effect_blocked" in image_result.quality_issues
    assert pdf_result.real_acceptance is False


def test_canonical_output_validates_against_v1_schema(tmp_path: Path) -> None:
    source = tmp_path / "sample.txt"
    source.write_text("traceable text", encoding="utf-8")
    document = ParserRouter().parse("txt", source, new_uuid7())
    repository_root = Path(__file__).resolve().parents[2]

    _assert_contract(document, repository_root)
