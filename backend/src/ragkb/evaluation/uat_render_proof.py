"""Independent local render representations for future UAT source proof."""

from __future__ import annotations

import csv
import hashlib
import json
from collections.abc import Mapping
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from openpyxl.utils.cell import range_boundaries
from pptx import Presentation
from pypdf import PdfReader


class RenderProofError(ValueError):
    pass


def _sha256_bytes(value: bytes) -> str:
    return hashlib.sha256(value, usedforsecurity=False).hexdigest()


def _canonical_hash(value: object) -> str:
    return _sha256_bytes(json.dumps(value, separators=(",", ":"), sort_keys=True).encode())


def _spreadsheet_text(path: Path, locator: Mapping[str, object]) -> str:
    if path.suffix.casefold() == ".csv":
        with path.open("r", encoding="utf-8", newline="") as handle:
            return "\n".join("\t".join(row) for row in csv.reader(handle))
    cell_range = locator.get("cell_range")
    sheet = locator.get("sheet")
    if not isinstance(cell_range, str) or not isinstance(sheet, str):
        raise RenderProofError("UAT_RENDER_PROOF_LOCATOR_UNSUPPORTED")
    workbook = load_workbook(path, read_only=True, data_only=True)
    if sheet not in workbook.sheetnames:
        raise RenderProofError("UAT_RENDER_PROOF_SHEET_MISSING")
    min_col, min_row, max_col, max_row = range_boundaries(cell_range)
    rows = workbook[sheet].iter_rows(
        min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col, values_only=True
    )
    return "\n".join(
        "\t".join("" if value is None else str(value) for value in row) for row in rows
    )


def independent_render_proof(
    *, category: str, source_path: Path, source_version_sha256: str, locator: Mapping[str, object]
) -> dict[str, object]:
    if _sha256_bytes(source_path.read_bytes()) != source_version_sha256:
        raise RenderProofError("UAT_RENDER_PROOF_SOURCE_HASH_MISMATCH")
    if category == "pdf_text":
        page = locator.get("page")
        if not isinstance(page, int) or page < 1:
            raise RenderProofError("UAT_RENDER_PROOF_LOCATOR_UNSUPPORTED")
        pages = PdfReader(source_path).pages
        if page > len(pages):
            raise RenderProofError("UAT_RENDER_PROOF_PAGE_MISSING")
        rendered_text = pages[page - 1].extract_text() or ""
    elif category == "pptx":
        slide = locator.get("slide")
        if not isinstance(slide, int) or slide < 1:
            raise RenderProofError("UAT_RENDER_PROOF_LOCATOR_UNSUPPORTED")
        slides = Presentation(str(source_path)).slides
        if slide > len(slides):
            raise RenderProofError("UAT_RENDER_PROOF_SLIDE_MISSING")
        rendered_text = "\n".join(
            shape.text
            for shape in slides[slide - 1].shapes
            if getattr(shape, "has_text_frame", False)
        )
    elif category == "spreadsheet":
        rendered_text = _spreadsheet_text(source_path, locator)
    elif category == "docx":
        document = Document(str(source_path))
        paragraphs = {paragraph._p: paragraph for paragraph in document.paragraphs}
        tables = {table._tbl: table for table in document.tables}
        parts = []
        for child in document.element.body.iterchildren():
            if child in paragraphs:
                parts.append(paragraphs[child].text)
            elif child in tables:
                parts.extend(
                    " | ".join(cell.text for cell in row.cells) for row in tables[child].rows
                )
        rendered_text = "\n".join(parts)
    else:
        raise RenderProofError("UAT_RENDER_PROOF_UNAVAILABLE")
    if not rendered_text.strip():
        raise RenderProofError("UAT_RENDER_PROOF_EMPTY")
    return {
        "revision": "uat-independent-render-proof:v1",
        "source_version_sha256": source_version_sha256,
        "locator_sha256": _canonical_hash(dict(locator)),
        "representation_sha256": hashlib.sha256(
            rendered_text.encode(), usedforsecurity=False
        ).hexdigest(),
        "rendered_text": rendered_text,
    }
