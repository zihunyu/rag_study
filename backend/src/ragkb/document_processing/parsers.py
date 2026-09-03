"""Native parsers and deferred routes behind the ParserPort contract."""

from __future__ import annotations

import csv
import hashlib
import wave
from collections.abc import Sequence
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

import openpyxl
import xlrd
from docx import Document as DocxDocument
from openpyxl.utils.cell import get_column_letter
from pptx import Presentation
from pypdf import PdfReader

from ragkb.contracts.ports import ParserPort, ParsingDeferred
from ragkb.domain.documents import CanonicalDocument, CanonicalNode, NodeType, SourceLocator
from ragkb.domain.ids import new_uuid7


def _checksum(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        while chunk := handle.read(1024 * 1024):
            digest.update(chunk)
    return digest.hexdigest()


def _document(
    path: Path,
    document_version_id: str,
    source_format: str,
    parser_revision: str,
    nodes: list[CanonicalNode],
    *,
    tables: list[dict[str, Any]] | None = None,
    quality_issues: list[str] | None = None,
) -> CanonicalDocument:
    return CanonicalDocument(
        document_version_id=document_version_id,
        language="und",
        source_format=source_format,
        nodes=tuple(nodes),
        parser_revision=parser_revision,
        normalization_revision="normalization:g1-v1",
        content_checksum=_checksum(path),
        tables=tuple(tables or []),
        quality_issues=tuple(quality_issues or []),
        real_acceptance=False,
    )


def _text_nodes(texts: list[str], *, locator_factory: Any) -> list[CanonicalNode]:
    nodes: list[CanonicalNode] = []
    offset = 0
    for index, raw in enumerate(texts):
        text = raw.strip()
        if not text:
            continue
        locator = locator_factory(index, offset, len(text))
        nodes.append(
            CanonicalNode(
                node_id=new_uuid7(),
                parent_node_id=None,
                node_type=NodeType.PARAGRAPH,
                original_text=text,
                display_text=text,
                locator=locator,
                metadata={"ordinal": index},
            )
        )
        offset += len(text) + 1
    return nodes


class PlainTextParser:
    revision = "plain-text-parser:g1-v1"

    def __init__(self, source_format: str) -> None:
        self.source_format = source_format

    def parse(self, source: Path, document_version_id: str) -> CanonicalDocument:
        text = source.read_text(encoding="utf-8")
        blocks = text.splitlines()
        nodes = _text_nodes(
            blocks,
            locator_factory=lambda _index, offset, length: SourceLocator(
                char_range=(offset, offset + length)
            ),
        )
        if not nodes:
            raise ParsingDeferred("PARSE_EMPTY", "text document has no usable content")
        return _document(source, document_version_id, self.source_format, self.revision, nodes)


class _HTMLExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.blocks: list[str] = []
        self._ignored_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag.casefold() in {"script", "style", "noscript"}:
            self._ignored_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag.casefold() in {"script", "style", "noscript"} and self._ignored_depth:
            self._ignored_depth -= 1

    def handle_data(self, data: str) -> None:
        if not self._ignored_depth and data.strip():
            self.blocks.append(data.strip())


class HTMLUploadParser:
    revision = "html-upload-parser:g1-v1"

    def parse(self, source: Path, document_version_id: str) -> CanonicalDocument:
        extractor = _HTMLExtractor()
        extractor.feed(source.read_text(encoding="utf-8"))
        nodes = _text_nodes(
            extractor.blocks,
            locator_factory=lambda _index, offset, length: SourceLocator(
                char_range=(offset, offset + length)
            ),
        )
        if not nodes:
            raise ParsingDeferred("PARSE_EMPTY", "HTML document has no visible text")
        return _document(source, document_version_id, "html", self.revision, nodes)


class TextPDFParser:
    revision = "pypdf-text:g1-v1"

    def parse(self, source: Path, document_version_id: str) -> CanonicalDocument:
        reader = PdfReader(str(source))
        nodes: list[CanonicalNode] = []
        for page_number, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if not text:
                continue
            nodes.append(
                CanonicalNode(
                    node_id=new_uuid7(),
                    parent_node_id=None,
                    node_type=NodeType.PARAGRAPH,
                    original_text=text,
                    display_text=text,
                    locator=SourceLocator(page=page_number),
                    metadata={},
                )
            )
        if not nodes:
            raise ParsingDeferred(
                "OCR_REQUIRED", "PDF has no usable text layer; MinerU/OCR is required"
            )
        return _document(source, document_version_id, "pdf_text", self.revision, nodes)


class ImageParserRoute:
    revision = "offline-ocr-stub:g4-v1"

    def __init__(self, source_format: str = "image") -> None:
        self.source_format = source_format

    def parse(self, source: Path, document_version_id: str) -> CanonicalDocument:
        if not source.read_bytes():
            raise ParsingDeferred("PARSE_EMPTY", "image fixture is empty")
        text = f"[offline OCR stub:{source.name}:{_checksum(source)[:12]}]"
        node = CanonicalNode(
            node_id=new_uuid7(),
            parent_node_id=None,
            node_type=NodeType.IMAGE,
            original_text=text,
            display_text=text,
            locator=SourceLocator(page=1, bbox=(0.0, 0.0, 1.0, 1.0)),
            metadata={"offline_stub": True},
        )
        return _document(
            source,
            document_version_id,
            self.source_format,
            self.revision,
            [node],
            quality_issues=["offline_ocr_stub_real_effect_blocked"],
        )


class OfflineOfficeConversionStubParser:
    revision = "offline-office-conversion-stub:g4-v1"

    def __init__(self, source_format: str) -> None:
        self.source_format = source_format

    def parse(self, source: Path, document_version_id: str) -> CanonicalDocument:
        if not source.read_bytes():
            raise ParsingDeferred("PARSE_EMPTY", "legacy Office fixture is empty")
        text = f"[offline Office conversion stub:{source.name}:{_checksum(source)[:12]}]"
        nodes = _text_nodes(
            [text],
            locator_factory=lambda _index, offset, length: SourceLocator(
                char_range=(offset, offset + length)
            ),
        )
        return _document(
            source,
            document_version_id,
            self.source_format,
            self.revision,
            nodes,
            quality_issues=["offline_office_conversion_stub_real_effect_blocked"],
        )


class OfflineASRStubParser:
    revision = "offline-asr-stub:g4-v1"

    def parse(self, source: Path, document_version_id: str) -> CanonicalDocument:
        if not source.read_bytes():
            raise ParsingDeferred("PARSE_EMPTY", "audio fixture is empty")
        duration = 1.0
        if source.suffix.casefold() == ".wav":
            try:
                with wave.open(str(source), "rb") as audio:
                    duration = audio.getnframes() / max(1, audio.getframerate())
            except wave.Error:
                duration = 1.0
        text = f"[offline ASR stub:{source.name}:{_checksum(source)[:12]}]"
        node = CanonicalNode(
            node_id=new_uuid7(),
            parent_node_id=None,
            node_type=NodeType.AUDIO,
            original_text=text,
            display_text=text,
            locator=SourceLocator(start_time=0.0, end_time=max(duration, 0.001)),
            metadata={"offline_stub": True},
        )
        return _document(
            source,
            document_version_id,
            "audio",
            self.revision,
            [node],
            quality_issues=["offline_asr_stub_real_effect_blocked"],
        )


class DOCXParser:
    revision = "python-docx:g1-v2"

    def parse(self, source: Path, document_version_id: str) -> CanonicalDocument:
        document = DocxDocument(str(source))
        paragraphs = {paragraph._p: paragraph for paragraph in document.paragraphs}
        tables = {table._tbl: table for table in document.tables}
        texts: list[str] = []
        for child in document.element.body.iterchildren():
            paragraph = paragraphs.get(child)
            if paragraph is not None:
                texts.append(paragraph.text)
                continue
            table = tables.get(child)
            if table is not None:
                texts.extend(" | ".join(cell.text for cell in row.cells) for row in table.rows)
        nodes = _text_nodes(
            texts,
            locator_factory=lambda _index, offset, length: SourceLocator(
                char_range=(offset, offset + length)
            ),
        )
        if not nodes:
            raise ParsingDeferred("PARSE_EMPTY", "DOCX document has no usable content")
        return _document(
            source,
            document_version_id,
            "docx",
            self.revision,
            nodes,
            quality_issues=["page_mapping_unavailable_in_native_docx_route"],
        )


class PPTXParser:
    revision = "python-pptx:g1-v2"

    def parse(self, source: Path, document_version_id: str) -> CanonicalDocument:
        presentation = Presentation(str(source))
        nodes: list[CanonicalNode] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            column_width = max(1, presentation.slide_width // 4)
            shapes = sorted(
                slide.shapes,
                key=lambda shape: (
                    shape.left // column_width,
                    shape.top,
                    shape.left,
                    shape.shape_id,
                ),
            )
            for shape in shapes:
                if getattr(shape, "has_table", False):
                    text = "\n".join(
                        " | ".join(cell.text for cell in row.cells) for row in shape.table.rows
                    ).strip()
                else:
                    text = str(getattr(shape, "text", "")).strip()
                if not text:
                    continue
                nodes.append(
                    CanonicalNode(
                        node_id=new_uuid7(),
                        parent_node_id=None,
                        node_type=NodeType.PARAGRAPH,
                        original_text=text,
                        display_text=text,
                        locator=SourceLocator(slide=slide_number),
                        metadata={},
                    )
                )
        if not nodes:
            raise ParsingDeferred("PARSE_EMPTY", "PPTX document has no usable content")
        return _document(source, document_version_id, "pptx", self.revision, nodes)


class SpreadsheetParser:
    revision = "spreadsheet-structure:g1-v2"

    @staticmethod
    def _node(sheet: str, row_number: int, values: Sequence[Any]) -> CanonicalNode | None:
        rendered = ["" if value is None else str(value) for value in values]
        if not any(value.strip() for value in rendered):
            return None
        text = " | ".join(rendered)
        end_column = get_column_letter(max(1, len(rendered)))
        return CanonicalNode(
            node_id=new_uuid7(),
            parent_node_id=None,
            node_type=NodeType.TABLE,
            original_text=text,
            display_text=text,
            locator=SourceLocator(
                sheet=sheet,
                cell_range=f"A{row_number}:{end_column}{row_number}",
                row=row_number,
            ),
            metadata={
                "row": row_number,
                "source_row_index": row_number,
                "column_count": len(rendered),
                "column_addresses": [
                    get_column_letter(index) for index in range(1, len(rendered) + 1)
                ],
                "values": rendered,
            },
        )

    def _xlsx(self, source: Path) -> tuple[list[CanonicalNode], list[dict[str, Any]]]:
        workbook = openpyxl.load_workbook(source, read_only=True, data_only=False)
        nodes: list[CanonicalNode] = []
        tables: list[dict[str, Any]] = []
        for sheet in workbook.worksheets:
            row_count = 0
            for row_number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                node = self._node(sheet.title, row_number, list(row))
                if node is not None:
                    nodes.append(node)
                    row_count += 1
            tables.append({"sheet": sheet.title, "non_empty_rows": row_count})
        workbook.close()
        return nodes, tables

    def _xls(self, source: Path) -> tuple[list[CanonicalNode], list[dict[str, Any]]]:
        workbook = xlrd.open_workbook(str(source), on_demand=True)
        nodes: list[CanonicalNode] = []
        tables: list[dict[str, Any]] = []
        for sheet in workbook.sheets():
            row_count = 0
            for row_index in range(sheet.nrows):
                node = self._node(sheet.name, row_index + 1, sheet.row_values(row_index))
                if node is not None:
                    nodes.append(node)
                    row_count += 1
            tables.append({"sheet": sheet.name, "non_empty_rows": row_count})
        workbook.release_resources()
        return nodes, tables

    def _csv(self, source: Path) -> tuple[list[CanonicalNode], list[dict[str, Any]]]:
        with source.open("r", encoding="utf-8-sig", newline="") as handle:
            rows = list(csv.reader(handle))
        nodes = [
            node
            for row_number, values in enumerate(rows, start=1)
            if (node := self._node("csv", row_number, values)) is not None
        ]
        return nodes, [{"sheet": "csv", "non_empty_rows": len(nodes)}]

    def parse(self, source: Path, document_version_id: str) -> CanonicalDocument:
        extension = source.suffix.casefold()
        if extension == ".xlsx":
            nodes, tables = self._xlsx(source)
            source_format = "xlsx"
        elif extension == ".xls":
            nodes, tables = self._xls(source)
            source_format = "xls"
        elif extension == ".csv":
            nodes, tables = self._csv(source)
            source_format = "csv"
        else:
            raise ParsingDeferred("PARSE_ROUTE_MISMATCH", "unsupported spreadsheet extension")
        if not nodes:
            raise ParsingDeferred("PARSE_EMPTY", "spreadsheet has no usable rows")
        return _document(
            source,
            document_version_id,
            source_format,
            self.revision,
            nodes,
            tables=tables,
        )


class ParserRouter:
    revision = "parser-router:g1-v1"

    def __init__(self) -> None:
        self._scanned_pdf_stub = ImageParserRoute("pdf_scanned")
        self._routes: dict[str, ParserPort] = {
            "txt": PlainTextParser("txt"),
            "markdown": PlainTextParser("markdown"),
            "html": HTMLUploadParser(),
            "pdf": TextPDFParser(),
            "image": ImageParserRoute(),
            "pdf_scanned": ImageParserRoute(),
            "doc": OfflineOfficeConversionStubParser("doc"),
            "ppt": OfflineOfficeConversionStubParser("ppt"),
            "docx": DOCXParser(),
            "pptx": PPTXParser(),
            "xlsx": SpreadsheetParser(),
            "xls": SpreadsheetParser(),
            "csv": SpreadsheetParser(),
            "audio": OfflineASRStubParser(),
        }

    def route(self, source_format: str) -> ParserPort:
        try:
            return self._routes[source_format]
        except KeyError as error:
            raise ParsingDeferred(
                "PARSE_ROUTE_UNAVAILABLE", f"no G1 parser route for {source_format}"
            ) from error

    def parse(
        self, source_format: str, source: Path, document_version_id: str
    ) -> CanonicalDocument:
        try:
            return self.route(source_format).parse(source, document_version_id)
        except ParsingDeferred as error:
            if source_format == "pdf" and error.code == "OCR_REQUIRED":
                return self._scanned_pdf_stub.parse(source, document_version_id)
            raise
