"""Structure-preserving DOCX, PPTX, XLSX, XLS, and CSV parsers."""

from __future__ import annotations

import csv
import re
from collections.abc import Sequence
from pathlib import Path
from typing import Any

import openpyxl
import xlrd
from docx import Document as DocxDocument
from openpyxl.utils.cell import get_column_letter
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from ragkb.contracts.ports import ParsingDeferred
from ragkb.document_processing.parser_common import canonical_document, text_nodes
from ragkb.domain.documents import CanonicalDocument, CanonicalNode, NodeType, SourceLocator
from ragkb.domain.ids import new_uuid7


class DOCXParser:
    revision = "python-docx:g1-v3"

    def parse(self, source: Path, document_version_id: str) -> CanonicalDocument:
        document = DocxDocument(str(source))
        paragraphs = {paragraph._p: paragraph for paragraph in document.paragraphs}
        tables = {table._tbl: table for table in document.tables}
        texts: list[str] = []
        node_types: list[NodeType] = []
        metadata: list[dict[str, Any]] = []
        for child in document.element.body.iterchildren():
            paragraph = paragraphs.get(child)
            if paragraph is not None:
                texts.append(paragraph.text)
                style_name = str(getattr(paragraph.style, "name", ""))
                heading = re.match(r"(?i)^heading\s+([1-6])$", style_name)
                node_types.append(NodeType.HEADING if heading else NodeType.PARAGRAPH)
                metadata.append({"heading_level": int(heading.group(1))} if heading else {})
                continue
            table = tables.get(child)
            if table is not None:
                rows = [" | ".join(cell.text for cell in row.cells) for row in table.rows]
                texts.extend(rows)
                node_types.extend(NodeType.TABLE for _ in rows)
                metadata.extend({"table_row": index + 1} for index in range(len(rows)))
        nodes = text_nodes(
            texts,
            locator_factory=lambda _index, offset, length: SourceLocator(
                char_range=(offset, offset + length)
            ),
            node_types=node_types,
            metadata=metadata,
        )
        if not nodes:
            raise ParsingDeferred("PARSE_EMPTY", "DOCX document has no usable content")
        return canonical_document(
            source,
            document_version_id,
            "docx",
            self.revision,
            nodes,
            quality_issues=["page_mapping_unavailable_in_native_docx_route"],
        )


class PPTXParser:
    revision = "python-pptx:g1-v3"

    @staticmethod
    def _node_type(shape: Any) -> NodeType:
        if getattr(shape, "is_placeholder", False) and shape.placeholder_format.type in {
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
        }:
            return NodeType.HEADING
        return NodeType.TABLE if getattr(shape, "has_table", False) else NodeType.PARAGRAPH

    def parse(self, source: Path, document_version_id: str) -> CanonicalDocument:
        presentation = Presentation(str(source))
        nodes: list[CanonicalNode] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            column_width = max(1, presentation.slide_width // 4)
            shapes = sorted(
                slide.shapes,
                key=lambda shape: (
                    shape.left // column_width,
                    shape.top,
                    shape.left,
                    shape.shape_id,
                ),
            )
            for shape in shapes:
                if getattr(shape, "has_table", False):
                    text = "\n".join(
                        " | ".join(cell.text for cell in row.cells) for row in shape.table.rows
                    ).strip()
                else:
                    text = str(getattr(shape, "text", "")).strip()
                if not text:
                    continue
                node_type = self._node_type(shape)
                nodes.append(
                    CanonicalNode(
                        node_id=new_uuid7(),
                        parent_node_id=None,
                        node_type=node_type,
                        original_text=text,
                        display_text=text,
                        locator=SourceLocator(slide=slide_number),
                        metadata={"heading_level": 1} if node_type is NodeType.HEADING else {},
                    )
                )
        if not nodes:
            raise ParsingDeferred("PARSE_EMPTY", "PPTX document has no usable content")
        return canonical_document(source, document_version_id, "pptx", self.revision, nodes)


class SpreadsheetParser:
    revision = "spreadsheet-structure:g1-v2"

    @staticmethod
    def _node(sheet: str, row_number: int, values: Sequence[Any]) -> CanonicalNode | None:
        rendered = ["" if value is None else str(value) for value in values]
        if not any(value.strip() for value in rendered):
            return None
        text = " | ".join(rendered)
        end_column = get_column_letter(max(1, len(rendered)))
        return CanonicalNode(
            node_id=new_uuid7(),
            parent_node_id=None,
            node_type=NodeType.TABLE,
            original_text=text,
            display_text=text,
            locator=SourceLocator(
                sheet=sheet,
                cell_range=f"A{row_number}:{end_column}{row_number}",
                row=row_number,
            ),
            metadata={
                "row": row_number,
                "source_row_index": row_number,
                "column_count": len(rendered),
                "column_addresses": [
                    get_column_letter(index) for index in range(1, len(rendered) + 1)
                ],
                "values": rendered,
            },
        )

    def _xlsx(self, source: Path) -> tuple[list[CanonicalNode], list[dict[str, Any]]]:
        workbook = openpyxl.load_workbook(source, read_only=True, data_only=False)
        nodes: list[CanonicalNode] = []
        tables: list[dict[str, Any]] = []
        for sheet in workbook.worksheets:
            row_count = 0
            for row_number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                node = self._node(sheet.title, row_number, list(row))
                if node is not None:
                    nodes.append(node)
                    row_count += 1
            tables.append({"sheet": sheet.title, "non_empty_rows": row_count})
        workbook.close()
        return nodes, tables

    def _xls(self, source: Path) -> tuple[list[CanonicalNode], list[dict[str, Any]]]:
        workbook = xlrd.open_workbook(str(source), on_demand=True)
        nodes: list[CanonicalNode] = []
        tables: list[dict[str, Any]] = []
        for sheet in workbook.sheets():
            row_count = 0
            for row_index in range(sheet.nrows):
                node = self._node(sheet.name, row_index + 1, sheet.row_values(row_index))
                if node is not None:
                    nodes.append(node)
                    row_count += 1
            tables.append({"sheet": sheet.name, "non_empty_rows": row_count})
        workbook.release_resources()
        return nodes, tables

    def _csv(self, source: Path) -> tuple[list[CanonicalNode], list[dict[str, Any]]]:
        with source.open("r", encoding="utf-8-sig", newline="") as handle:
            rows = list(csv.reader(handle))
        nodes = [
            node
            for row_number, values in enumerate(rows, start=1)
            if (node := self._node("csv", row_number, values)) is not None
        ]
        return nodes, [{"sheet": "csv", "non_empty_rows": len(nodes)}]

    def parse(self, source: Path, document_version_id: str) -> CanonicalDocument:
        extension = source.suffix.casefold()
        if extension == ".xlsx":
            nodes, tables = self._xlsx(source)
            source_format = "xlsx"
        elif extension == ".xls":
            nodes, tables = self._xls(source)
            source_format = "xls"
        elif extension == ".csv":
            nodes, tables = self._csv(source)
            source_format = "csv"
        else:
            raise ParsingDeferred("PARSE_ROUTE_MISMATCH", "unsupported spreadsheet extension")
        if not nodes:
            raise ParsingDeferred("PARSE_EMPTY", "spreadsheet has no usable rows")
        return canonical_document(
            source,
            document_version_id,
            source_format,
            self.revision,
            nodes,
            tables=tables,
        )
